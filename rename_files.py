#!/usr/bin/env python3
"""
Rename files intelligently using Ollama vision language models.
Analyzes file content and generates meaningful names.
"""

from __future__ import annotations

import argparse
import base64
import json
import re
import tempfile
from pathlib import Path
from typing import Iterable

import fitz  # PyMuPDF
import requests


def encode_image_to_base64(image_path: Path) -> str:
    """Encode an image file to base64 string."""
    return base64.b64encode(image_path.read_bytes()).decode("utf-8")


def call_ollama_vision(
    model: str,
    prompt: str,
    images: Iterable[Path],
    ollama_url: str = "http://127.0.0.1:11434",
) -> str:
    """Call Ollama vision model with images and a prompt."""
    payload = {
        "model": model,
        "prompt": prompt,
        "images": [encode_image_to_base64(path) for path in images],
        "stream": False,
    }

    response = requests.post(f"{ollama_url}/api/generate", json=payload, timeout=600)
    response.raise_for_status()
    data = response.json()
    return str(data.get("response", "")).strip()


def call_ollama_text(
    model: str,
    prompt: str,
    ollama_url: str = "http://127.0.0.1:11434",
) -> str:
    """Call Ollama model for text-only analysis."""
    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
    }

    response = requests.post(f"{ollama_url}/api/generate", json=payload, timeout=600)
    response.raise_for_status()
    data = response.json()
    return str(data.get("response", "")).strip()


def extract_pdf_first_page(pdf_path: Path, dpi: int = 150) -> Path:
    """Extract first page of PDF as an image."""
    doc = fitz.open(pdf_path)
    page = doc.load_page(0)
    
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=matrix)
    
    temp_dir = Path(tempfile.gettempdir()) / "smartname"
    temp_dir.mkdir(exist_ok=True)
    image_path = temp_dir / f"{pdf_path.stem}_page1.png"
    pix.save(image_path)
    
    doc.close()
    return image_path


def extract_video_frame(video_path: Path) -> Path | None:
    """Extract a frame from video using ffmpeg if available."""
    try:
        import subprocess
        
        temp_dir = Path(tempfile.gettempdir()) / "smartname"
        temp_dir.mkdir(exist_ok=True)
        image_path = temp_dir / f"{video_path.stem}_frame.jpg"
        
        # Extract frame at 1 second
        result = subprocess.run(
            [
                "ffmpeg", "-i", str(video_path),
                "-ss", "00:00:01",
                "-vframes", "1",
                "-y",
                str(image_path)
            ],
            capture_output=True,
            text=True
        )
        
        if result.returncode == 0 and image_path.exists():
            return image_path
    except (ImportError, FileNotFoundError):
        pass
    
    return None


def read_text_snippet(file_path: Path, max_chars: int = 2000) -> str:
    """Read the beginning of a text file."""
    try:
        content = file_path.read_text(encoding="utf-8")
        if len(content) > max_chars:
            return content[:max_chars] + "..."
        return content
    except Exception:
        return ""


def extract_docx_content(docx_path: Path, max_chars: int = 2000) -> str:
    """Extract text content from a DOCX file."""
    try:
        from docx import Document
        doc = Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)
        if len(content) > max_chars:
            return content[:max_chars] + "..."
        return content
    except ImportError:
        return ""
    except Exception:
        return ""


def extract_pptx_content(pptx_path: Path, max_chars: int = 2000) -> str:
    """Extract text content from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        texts = []
        
        # Extract text from first few slides
        for slide in prs.slides[:5]:  # First 5 slides
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
        
        content = "\n".join(texts)
        if len(content) > max_chars:
            return content[:max_chars] + "..."
        return content
    except ImportError:
        return ""
    except Exception:
        return ""


def docx_to_images(docx_path: Path, dpi: int = 150) -> list[Path]:
    """
    Convert first page of DOCX to image using LibreOffice if available.
    Falls back to text extraction.
    """
    try:
        import subprocess
        
        temp_dir = Path(tempfile.gettempdir()) / "smartname"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert to PDF first using LibreOffice
        pdf_path = temp_dir / f"{docx_path.stem}.pdf"
        result = subprocess.run(
            [
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(temp_dir), str(docx_path)
            ],
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode == 0 and pdf_path.exists():
            # Convert PDF to image
            image_path = extract_pdf_first_page(pdf_path, dpi)
            return [image_path]
    except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
        pass
    
    return []


def pptx_to_images(pptx_path: Path, dpi: int = 150) -> list[Path]:
    """
    Convert first slide of PPTX to image using LibreOffice if available.
    Falls back to text extraction.
    """
    try:
        import subprocess
        
        temp_dir = Path(tempfile.gettempdir()) / "smartname"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert to PDF first using LibreOffice
        pdf_path = temp_dir / f"{pptx_path.stem}.pdf"
        result = subprocess.run(
            [
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(temp_dir), str(pptx_path)
            ],
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode == 0 and pdf_path.exists():
            # Convert PDF to image
            image_path = extract_pdf_first_page(pdf_path, dpi)
            return [image_path]
    except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
        pass
    
    return []


def apply_casing(text: str, case_style: str) -> str:
    """
    Apply a casing style to text.
    
    Styles:
    - snake: snake_case (default)
    - kebab: kebab-case
    - camel: camelCase
    - pascal: PascalCase
    - lower: lowercase with spaces
    - title: Title Case With Spaces
    """
    # First normalize to words
    words = re.sub(r'[_\-\s]+', ' ', text).strip().split()
    
    if case_style == "snake":
        return "_".join(w.lower() for w in words)
    elif case_style == "kebab":
        return "-".join(w.lower() for w in words)
    elif case_style == "camel":
        if not words:
            return ""
        return words[0].lower() + "".join(w.capitalize() for w in words[1:])
    elif case_style == "pascal":
        return "".join(w.capitalize() for w in words)
    elif case_style == "lower":
        return " ".join(w.lower() for w in words)
    elif case_style == "title":
        return " ".join(w.capitalize() for w in words)
    else:
        return "_".join(w.lower() for w in words)  # default to snake_case


def sanitize_filename(name: str, max_length: int = 100, case_style: str = "snake") -> str:
    """
    Clean up a suggested filename to be filesystem-safe.
    Remove/replace invalid characters and limit length.
    """
    # Remove or replace invalid characters
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    # Replace multiple spaces/underscores with single ones
    name = re.sub(r'[\s_]+', '_', name)
    # Remove leading/trailing spaces and underscores
    name = name.strip(' _.')
    
    # Apply casing style
    name = apply_casing(name, case_style)
    
    # Limit length
    if len(name) > max_length:
        # Try to cut at word boundary
        separator = "_" if case_style in ["snake", "lower", "title"] else "-" if case_style == "kebab" else ""
        if separator and separator in name:
            name = name[:max_length].rsplit(separator, 1)[0]
        else:
            name = name[:max_length]
    
    return name if name else "unnamed"


# Default categories for file organization
DEFAULT_CATEGORIES = [
    "books",
    "photos",
    "figures",
    "documents",
    "presentations",
    "screenshots",
    "art",
    "code",
    "other",
]


def categorize_file(
    file_path: Path,
    model: str,
    ollama_url: str,
    categories: list[str] | None = None,
    dpi: int = 150,
) -> str:
    """
    Analyze a file and categorize it into one of the predefined categories.
    Returns the category name.
    """
    if categories is None:
        categories = DEFAULT_CATEGORIES
    
    file_ext = file_path.suffix.lower()
    category_list = ", ".join(categories)
    
    print(f"\nCategorizing: {file_path.name}")
    
    # Handle images
    if file_ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        prompt = (
            f"Analyze this image and categorize it into ONE of these categories: {category_list}.\n"
            "Consider the content and purpose of the image.\n"
            "Examples:\n"
            "- 'books': textbook pages, book covers, scanned book pages\n"
            "- 'photos': personal photos, portraits, landscapes, snapshots\n"
            "- 'figures': charts, graphs, diagrams, scientific figures, infographics\n"
            "- 'documents': forms, letters, receipts, official documents\n"
            "- 'screenshots': screen captures, UI screenshots, app interfaces\n"
            "- 'art': artwork, illustrations, paintings, creative images\n"
            "- 'presentations': presentation slides, slide decks\n"
            "- 'code': code snippets, programming screenshots\n"
            "- 'other': anything that doesn't fit the above categories\n\n"
            "Respond with ONLY the category name, nothing else."
        )
        category = call_ollama_vision(model, prompt, [file_path], ollama_url)
    
    # Handle PDFs
    elif file_ext == ".pdf":
        temp_image = extract_pdf_first_page(file_path, dpi)
        prompt = (
            f"Analyze this PDF document and categorize it into ONE of these categories: {category_list}.\n"
            "Consider the content and purpose of the document.\n"
            "Examples:\n"
            "- 'books': textbooks, ebooks, book chapters\n"
            "- 'documents': forms, letters, reports, official documents\n"
            "- 'presentations': presentation slides, slide decks\n"
            "- 'figures': research papers with figures, scientific documents\n"
            "- 'other': anything that doesn't fit the above categories\n\n"
            "Respond with ONLY the category name, nothing else."
        )
        category = call_ollama_vision(model, prompt, [temp_image], ollama_url)
    
    # Handle DOCX files
    elif file_ext == ".docx":
        # Try visual conversion first
        images = docx_to_images(file_path, dpi)
        if images:
            prompt = (
                f"Analyze this Word document and categorize it into ONE of these categories: {category_list}.\n"
                "Consider the content and purpose.\n"
                "Respond with ONLY the category name, nothing else."
            )
            category = call_ollama_vision(model, prompt, images, ollama_url)
        else:
            # Fall back to text extraction
            content = extract_docx_content(file_path)
            if not content:
                print("  Could not extract DOCX content, defaulting to 'documents'")
                return "documents"
            
            prompt = (
                f"This is the content of a Word document:\n\n{content}\n\n"
                f"Categorize it into ONE of these categories: {category_list}.\n"
                "Respond with ONLY the category name, nothing else."
            )
            category = call_ollama_text(model, prompt, ollama_url)
    
    # Handle PPTX files
    elif file_ext == ".pptx":
        # Try visual conversion first
        images = pptx_to_images(file_path, dpi)
        if images:
            prompt = (
                f"Analyze this PowerPoint presentation and categorize it into ONE of these categories: {category_list}.\n"
                "Respond with ONLY the category name, nothing else."
            )
            category = call_ollama_vision(model, prompt, images, ollama_url)
        else:
            # Fall back to text extraction
            content = extract_pptx_content(file_path)
            if not content:
                print("  Could not extract PPTX content, defaulting to 'presentations'")
                return "presentations"
            
            prompt = (
                f"This is the content of a PowerPoint presentation:\n\n{content}\n\n"
                f"Categorize it into ONE of these categories: {category_list}.\n"
                "Respond with ONLY the category name, nothing else."
            )
            category = call_ollama_text(model, prompt, ollama_url)
    
    # Handle videos
    elif file_ext in {".mov", ".mp4", ".avi", ".mkv", ".webm"}:
        frame_path = extract_video_frame(file_path)
        if frame_path:
            prompt = (
                f"Analyze this video frame and categorize the video into ONE of these categories: {category_list}.\n"
                "Respond with ONLY the category name, nothing else."
            )
            category = call_ollama_vision(model, prompt, [frame_path], ollama_url)
        else:
            print("  Could not extract video frame, defaulting to 'other'")
            return "other"
    
    # Handle text files and notebooks
    elif file_ext in {".txt", ".md", ".ipynb", ".py", ".js", ".json"}:
        content = read_text_snippet(file_path)
        if not content:
            print("  Could not read file content, defaulting to 'code'")
            return "code"
        
        prompt = (
            f"This is the content of a {file_ext} file:\n\n{content}\n\n"
            f"Categorize it into ONE of these categories: {category_list}.\n"
            "Respond with ONLY the category name, nothing else."
        )
        category = call_ollama_text(model, prompt, ollama_url)
    
    else:
        print(f"  Unsupported file type: {file_ext}, defaulting to 'other'")
        return "other"
    
    # Clean up the response
    category = category.strip().strip('"\'`.,!?;:').lower()
    
    # Validate category - try exact match first
    if category in categories:
        print(f"  Category: {category}")
        return category
    
    # Try to find partial match (e.g., "art." -> "art")
    for valid_cat in categories:
        if category.startswith(valid_cat) or valid_cat in category:
            print(f"  Category: {valid_cat} (matched from '{category}')")
            return valid_cat
    
    # Default to 'other' if no match found
    print(f"  VLM returned invalid category '{category}', defaulting to 'other'")
    return "other"


def generate_filename_for_file(
    file_path: Path,
    model: str,
    ollama_url: str,
    dpi: int = 150,
    case_style: str = "snake",
) -> str:
    """
    Analyze a file and generate a descriptive filename.
    Returns the suggested name without extension.
    """
    file_ext = file_path.suffix.lower()
    
    print(f"\nAnalyzing: {file_path.name}")
    
    # Handle images
    if file_ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        prompt = (
            "Analyze this image and suggest a concise, descriptive filename "
            "(5-8 words max, use underscores instead of spaces). "
            "Focus on the main subject or content. "
            "Only respond with the filename, nothing else."
        )
        suggested = call_ollama_vision(model, prompt, [file_path], ollama_url)
    
    # Handle PDFs
    elif file_ext == ".pdf":
        temp_image = extract_pdf_first_page(file_path, dpi)
        prompt = (
            "This is the first page of a PDF document. "
            "Analyze it and suggest a concise, descriptive filename "
            "(5-8 words max, use underscores instead of spaces). "
            "Focus on the document's topic or title. "
            "Only respond with the filename, nothing else."
        )
        suggested = call_ollama_vision(model, prompt, [temp_image], ollama_url)
    
    # Handle DOCX files
    elif file_ext == ".docx":
        # Try visual conversion first
        images = docx_to_images(file_path, dpi)
        if images:
            prompt = (
                "This is the first page of a Word document. "
                "Analyze it and suggest a concise, descriptive filename "
                "(5-8 words max, use underscores instead of spaces). "
                "Focus on the document's topic or title. "
                "Only respond with the filename, nothing else."
            )
            suggested = call_ollama_vision(model, prompt, images, ollama_url)
        else:
            # Fall back to text extraction
            content = extract_docx_content(file_path)
            if not content:
                print("  Could not extract DOCX content")
                return file_path.stem
            
            prompt = (
                f"This is the content of a Word document:\n\n"
                f"{content}\n\n"
                "Suggest a concise, descriptive filename based on the content "
                "(5-8 words max, use underscores instead of spaces). "
                "Only respond with the filename, nothing else."
            )
            suggested = call_ollama_text(model, prompt, ollama_url)
    
    # Handle PPTX files
    elif file_ext == ".pptx":
        # Try visual conversion first
        images = pptx_to_images(file_path, dpi)
        if images:
            prompt = (
                "This is the first slide of a PowerPoint presentation. "
                "Analyze it and suggest a concise, descriptive filename "
                "(5-8 words max, use underscores instead of spaces). "
                "Focus on the presentation's topic or title. "
                "Only respond with the filename, nothing else."
            )
            suggested = call_ollama_vision(model, prompt, images, ollama_url)
        else:
            # Fall back to text extraction
            content = extract_pptx_content(file_path)
            if not content:
                print("  Could not extract PPTX content")
                return file_path.stem
            
            prompt = (
                f"This is the content of a PowerPoint presentation:\n\n"
                f"{content}\n\n"
                "Suggest a concise, descriptive filename based on the content "
                "(5-8 words max, use underscores instead of spaces). "
                "Only respond with the filename, nothing else."
            )
            suggested = call_ollama_text(model, prompt, ollama_url)
    
    # Handle videos
    elif file_ext in {".mov", ".mp4", ".avi", ".mkv", ".webm"}:
        frame_path = extract_video_frame(file_path)
        if frame_path:
            prompt = (
                "This is a frame from a video. "
                "Analyze it and suggest a concise, descriptive filename for the video "
                "(5-8 words max, use underscores instead of spaces). "
                "Only respond with the filename, nothing else."
            )
            suggested = call_ollama_vision(model, prompt, [frame_path], ollama_url)
        else:
            print("  Could not extract video frame (ffmpeg not available)")
            return file_path.stem
    
    # Handle text files and notebooks
    elif file_ext in {".txt", ".md", ".ipynb", ".py", ".js", ".json"}:
        content = read_text_snippet(file_path)
        if not content:
            print("  Could not read file content")
            return file_path.stem
        
        prompt = (
            f"This is the content of a {file_ext} file:\n\n"
            f"{content}\n\n"
            "Suggest a concise, descriptive filename based on the content "
            "(5-8 words max, use underscores instead of spaces). "
            "Only respond with the filename, nothing else."
        )
        suggested = call_ollama_text(model, prompt, ollama_url)
    
    else:
        print(f"  Unsupported file type: {file_ext}")
        return file_path.stem
    
    # Clean up the suggestion
    suggested = suggested.strip().strip('"\'`')
    # Remove any file extensions that the model might have added
    suggested = re.sub(r'\.(txt|pdf|png|jpg|jpeg|mov|mp4|ipynb|py|js|json|md|docx|pptx)$', '', suggested, flags=re.IGNORECASE)
    suggested = sanitize_filename(suggested, case_style=case_style)
    
    print(f"  Suggested: {suggested}{file_ext}")
    return suggested


def rename_files_in_directory(
    directory: Path,
    model: str,
    ollama_url: str,
    dry_run: bool = True,
    dpi: int = 150,
    case_style: str = "snake",
) -> None:
    """
    Rename all supported files in a directory based on their content.
    """
    supported_extensions = {
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp",  # Images
        ".pdf",  # PDFs
        ".docx", ".pptx",  # Office documents
        ".mov", ".mp4", ".avi", ".mkv", ".webm",  # Videos
        ".txt", ".md", ".ipynb", ".py", ".js", ".json",  # Text/Code
    }
    
    files = [
        f for f in directory.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]
    
    if not files:
        print(f"No supported files found in {directory}")
        return
    
    print(f"Found {len(files)} files to process")
    print(f"Model: {model}")
    print(f"Case style: {case_style}")
    print(f"Mode: {'DRY RUN' if dry_run else 'RENAMING'}")
    
    renames = []
    
    for file_path in files:
        try:
            new_name = generate_filename_for_file(file_path, model, ollama_url, dpi, case_style)
            new_path = file_path.parent / f"{new_name}{file_path.suffix}"
            
            # Avoid overwriting existing files
            counter = 1
            while new_path.exists() and new_path != file_path:
                new_path = file_path.parent / f"{new_name}_{counter}{file_path.suffix}"
                counter += 1
            
            if new_path != file_path:
                renames.append((file_path, new_path))
        
        except Exception as e:
            print(f"  Error processing {file_path.name}: {e}")
    
    # Execute renames
    print("\n" + "=" * 60)
    print("RENAME SUMMARY")
    print("=" * 60)
    
    for old_path, new_path in renames:
        print(f"\n{old_path.name}")
        print(f"  → {new_path.name}")
        
        if not dry_run:
            try:
                old_path.rename(new_path)
                print("  ✓ Renamed")
            except Exception as e:
                print(f"  ✗ Error: {e}")
    
    if dry_run:
        print("\n" + "=" * 60)
        print("This was a DRY RUN. Use --execute to actually rename files.")
        print("=" * 60)


def organize_files_in_directory(
    directory: Path,
    model: str,
    ollama_url: str,
    dry_run: bool = True,
    dpi: int = 150,
    case_style: str = "snake",
    categories: list[str] | None = None,
    rename_files: bool = False,
) -> None:
    """
    Organize files into categorized folders based on their content.
    Optionally rename files within folders.
    """
    if categories is None:
        categories = DEFAULT_CATEGORIES
    
    supported_extensions = {
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp",  # Images
        ".pdf",  # PDFs
        ".docx", ".pptx",  # Office documents
        ".mov", ".mp4", ".avi", ".mkv", ".webm",  # Videos
        ".txt", ".md", ".ipynb", ".py", ".js", ".json",  # Text/Code
    }
    
    files = [
        f for f in directory.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]
    
    if not files:
        print(f"No supported files found in {directory}")
        return
    
    print(f"Found {len(files)} files to organize")
    print(f"Model: {model}")
    print(f"Categories: {', '.join(categories)}")
    print(f"Rename files: {'Yes' if rename_files else 'No'}")
    print(f"Case style: {case_style}")
    print(f"Mode: {'DRY RUN' if dry_run else 'ORGANIZING'}")
    
    # Categorize all files
    file_categories: dict[Path, str] = {}
    
    for file_path in files:
        try:
            category = categorize_file(file_path, model, ollama_url, categories, dpi)
            file_categories[file_path] = category
        except Exception as e:
            print(f"  Error categorizing {file_path.name}: {e}")
            file_categories[file_path] = "other"
    
    # Organize files by category
    moves: list[tuple[Path, Path]] = []
    
    for file_path, category in file_categories.items():
        # Create category folder path
        category_folder = directory / category
        
        # Determine new filename
        if rename_files:
            try:
                new_name = generate_filename_for_file(file_path, model, ollama_url, dpi, case_style)
                new_path = category_folder / f"{new_name}{file_path.suffix}"
            except Exception as e:
                print(f"  Error generating name for {file_path.name}: {e}, using original name")
                new_path = category_folder / file_path.name
        else:
            new_path = category_folder / file_path.name
        
        # Avoid overwriting existing files
        counter = 1
        original_new_path = new_path
        while new_path.exists() and new_path != file_path:
            stem = original_new_path.stem
            new_path = category_folder / f"{stem}_{counter}{file_path.suffix}"
            counter += 1
        
        if new_path != file_path:
            moves.append((file_path, new_path))
    
    # Execute organization
    print("\n" + "=" * 60)
    print("ORGANIZATION SUMMARY")
    print("=" * 60)
    
    # Group by category for better display
    by_category: dict[str, list[tuple[Path, Path]]] = {}
    for old_path, new_path in moves:
        category = file_categories[old_path]
        if category not in by_category:
            by_category[category] = []
        by_category[category].append((old_path, new_path))
    
    for category in sorted(by_category.keys()):
        print(f"\n[{category.upper()}]")
        for old_path, new_path in by_category[category]:
            print(f"  {old_path.name}")
            print(f"    → {new_path.relative_to(directory)}")
            
            if not dry_run:
                try:
                    # Create category folder if it doesn't exist
                    new_path.parent.mkdir(parents=True, exist_ok=True)
                    # Move the file
                    old_path.rename(new_path)
                    print("    ✓ Organized")
                except Exception as e:
                    print(f"    ✗ Error: {e}")
    
    if dry_run:
        print("\n" + "=" * 60)
        print("This was a DRY RUN. Use --execute to actually organize files.")
        print("=" * 60)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Rename and organize files intelligently using Ollama vision models"
    )
    parser.add_argument(
        "directory",
        type=Path,
        help="Directory containing files to rename/organize",
    )
    parser.add_argument(
        "--model",
        default="llava:latest",
        help="Ollama model to use (default: llava:latest)",
    )
    parser.add_argument(
        "--ollama-url",
        default="http://127.0.0.1:11434",
        help="Ollama server URL (default: http://127.0.0.1:11434)",
    )
    parser.add_argument(
        "--execute",
        action="store_true",
        help="Actually rename/organize files (default is dry-run)",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=150,
        help="DPI for PDF rendering (default: 150)",
    )
    parser.add_argument(
        "--case",
        dest="case_style",
        choices=["snake", "kebab", "camel", "pascal", "lower", "title"],
        default="snake",
        help="Casing style for filenames (default: snake). Options: "
             "snake (snake_case), kebab (kebab-case), camel (camelCase), "
             "pascal (PascalCase), lower (lowercase with spaces), "
             "title (Title Case With Spaces)",
    )
    parser.add_argument(
        "--organize",
        action="store_true",
        help="Organize files into categorized folders instead of just renaming",
    )
    parser.add_argument(
        "--rename-in-folders",
        action="store_true",
        help="When organizing, also rename files within folders (requires --organize)",
    )
    parser.add_argument(
        "--categories",
        nargs="+",
        help="Custom categories for organization (default: books, photos, figures, "
             "documents, presentations, screenshots, art, code, other)",
    )
    
    args = parser.parse_args()
    
    if not args.directory.is_dir():
        print(f"Error: {args.directory} is not a directory")
        return
    
    if args.rename_in_folders and not args.organize:
        print("Warning: --rename-in-folders requires --organize, ignoring flag")
        args.rename_in_folders = False
    
    # Choose mode: organize or rename
    if args.organize:
        organize_files_in_directory(
            directory=args.directory,
            model=args.model,
            ollama_url=args.ollama_url,
            dry_run=not args.execute,
            dpi=args.dpi,
            case_style=args.case_style,
            categories=args.categories,
            rename_files=args.rename_in_folders,
        )
    else:
        rename_files_in_directory(
            directory=args.directory,
            model=args.model,
            ollama_url=args.ollama_url,
            dry_run=not args.execute,
            dpi=args.dpi,
            case_style=args.case_style,
        )


if __name__ == "__main__":
    main()

